import json
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.data import ChartData
import pandas as pd
import os
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def get_theme_color(scheme_color):
    """Convert scheme color string to MSO_THEME_COLOR enum."""
    try:
        if not scheme_color or not isinstance(scheme_color, str):
            logger.warning(f"Invalid scheme color value: {scheme_color}")
            return MSO_THEME_COLOR.ACCENT_1
            
        if not scheme_color.startswith('SCHEME'):
            logger.warning(f"Not a scheme color: {scheme_color}")
            return MSO_THEME_COLOR.ACCENT_1
            
        try:
            number = int(scheme_color.split('(')[1].strip(')'))
        except (IndexError, ValueError) as e:
            logger.warning(f"Could not parse scheme color number from {scheme_color}: {e}")
            return MSO_THEME_COLOR.ACCENT_1

        theme_colors = {
            1: MSO_THEME_COLOR.ACCENT_1,
            2: MSO_THEME_COLOR.ACCENT_2,
            3: MSO_THEME_COLOR.ACCENT_3,
            4: MSO_THEME_COLOR.ACCENT_4,
            5: MSO_THEME_COLOR.ACCENT_5,
            6: MSO_THEME_COLOR.ACCENT_6,
            7: MSO_THEME_COLOR.ACCENT_7,
            8: MSO_THEME_COLOR.ACCENT_8,
            9: MSO_THEME_COLOR.ACCENT_9,
            10: MSO_THEME_COLOR.ACCENT_10,
            11: MSO_THEME_COLOR.ACCENT_11,
            12: MSO_THEME_COLOR.ACCENT_12,
            13: MSO_THEME_COLOR.BACKGROUND_1,
            14: MSO_THEME_COLOR.BACKGROUND_2,
            15: MSO_THEME_COLOR.TEXT_1,
            16: MSO_THEME_COLOR.TEXT_2,
            17: MSO_THEME_COLOR.TEXT_3,
            18: MSO_THEME_COLOR.TEXT_4,
            19: MSO_THEME_COLOR.TEXT_5,
            20: MSO_THEME_COLOR.TEXT_6,
            21: MSO_THEME_COLOR.TEXT_7,
            22: MSO_THEME_COLOR.TEXT_8,
            23: MSO_THEME_COLOR.TEXT_9,
            24: MSO_THEME_COLOR.TEXT_10,
        }
        
        if number not in theme_colors:
            logger.warning(f"Unknown theme color number {number}, defaulting to ACCENT_1")
            return MSO_THEME_COLOR.ACCENT_1
            
        return theme_colors[number]
    except Exception as e:
        logger.warning(f"Error processing theme color {scheme_color}: {e}")
        return MSO_THEME_COLOR.ACCENT_1

def update_text_blocks(prs, jsonl_path):
    """Update text blocks in the presentation."""
    logger.info("Updating text blocks...")
    translations = {}
    with open(jsonl_path, 'r', encoding='utf-8') as f:
        for line in f:
            block = json.loads(line)
            key = (
                block['slide_index'],
                block['shape_index'],
                block['paragraph_index'],
                block['run_index']
            )
            translations[key] = block.get('french_text', block.get('text', ''))

    updates = 0
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_index, run in enumerate(paragraph.runs):
                    key = (slide_index, shape_index, paragraph_index, run_index)
                    if key in translations:
                        run.text = translations[key]
                        updates += 1
    
    logger.info(f"Updated {updates} text blocks")

def update_tables(prs, jsonl_path):
    """Update tables in the presentation."""
    logger.info("Updating tables...")
    updates = 0
    with open(jsonl_path, 'r', encoding='utf-8') as f:
        for line in f:
            try:
                record = json.loads(line.strip())
                slide_idx = record['slide_index']
                shape_idx = record['shape_index']
                
                try:
                    slide = prs.slides[slide_idx]
                    shape = slide.shapes[shape_idx]
                    if not shape.has_table:
                        continue
                    table = shape.table
                except Exception as e:
                    logger.error(f"Error accessing table on slide {slide_idx}: {e}")
                    continue
                
                data = record['data']
                if not data:
                    continue
                
                for row_idx, row_data in enumerate(data):
                    if row_idx >= len(table.rows):
                        continue
                        
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx >= len(table.columns):
                            continue
                            
                        try:
                            cell = table.cell(row_idx, col_idx)
                            
                            # Clear existing content
                            for paragraph in cell.text_frame.paragraphs:
                                p = paragraph._element
                                p.getparent().remove(p)
                            
                            # Add new content
                            if isinstance(cell_data, dict):
                                for para in cell_data.get('paragraphs', []):
                                    p = cell.text_frame.add_paragraph()
                                    if para.get('is_bullet'):
                                        p.level = para.get('bullet_level', 0)
                                    for run_data in para.get('runs', []):
                                        r = p.add_run()
                                        if run_data.get('french_text'):
                                            r.text = run_data['french_text']
                                        else:
                                            r.text = run_data.get('text', '')
                                        
                                        # Apply formatting
                                        if run_data.get('bold'):
                                            r.font.bold = True
                                        if run_data.get('italic'):
                                            r.font.italic = True
                                        if run_data.get('size'):
                                            try:
                                                r.font.size = Pt(float(run_data['size']))
                                            except:
                                                pass
                                        if run_data.get('font'):
                                            r.font.name = run_data['font']
                                        if run_data.get('color'):
                                            try:
                                                if run_data['color'].startswith('SCHEME'):
                                                    theme_color = get_theme_color(run_data['color'])
                                                    r.font.color.theme_color = theme_color
                                                else:
                                                    color = run_data['color'].lstrip('#')
                                                    r.font.color.rgb = RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))
                                            except Exception as e:
                                                logger.warning(f"Could not apply color formatting: {e}")
                                        if run_data.get('url'):
                                            r.hyperlink.address = run_data['url']
                                        updates += 1
                            else:
                                p = cell.text_frame.add_paragraph()
                                r = p.add_run()
                                r.text = str(cell_data)
                                updates += 1
                                
                        except Exception as e:
                            logger.error(f"Error updating cell at row {row_idx}, col {col_idx} on slide {slide_idx}: {e}")
                            continue
                            
            except Exception as e:
                logger.error(f"Error processing table record: {e}")
                continue
    
    logger.info(f"Updated {updates} table cells")

def update_charts(prs, jsonl_path):
    """Update charts in the presentation."""
    logger.info("Updating charts...")
    translations = {}
    with open(jsonl_path, 'r', encoding='utf-8') as f:
        for line in f:
            data = json.loads(line)
            if 'text' in data and 'french_text' in data:
                key = (data.get('slide_index'), data.get('shape_index'), data['text'])
                translations[key] = data['french_text']

    updates = 0
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not isinstance(shape, GraphicFrame):
                continue
                
            try:
                chart = shape.chart
                
                # Update chart title
                if chart.has_title:
                    title = chart.chart_title.text_frame.text.strip()
                    title_key = (slide_index, shape_index, title)
                    if title_key in translations:
                        chart.chart_title.text_frame.text = translations[title_key]
                        updates += 1
                
                # Update chart data for pie/doughnut charts
                if chart.chart_type in [-4120, -4121]:  # DOUGHNUT = -4120, PIE = -4121
                    chart_data_bytes = chart.part.chart_workbook.xlsx_part.blob
                    excel_path = f'temp_chart_data_{slide_index}_{shape_index}.xlsx'
                    
                    with open(excel_path, 'wb') as f:
                        f.write(chart_data_bytes)
                    
                    df = pd.read_excel(excel_path)
                    categories = []
                    values = []
                    
                    for idx, row in df.iterrows():
                        category_name = None
                        if isinstance(row.name, str):
                            category_name = row.name.strip()
                        if not category_name:
                            row_str = str(row)
                            parts = row_str.split()
                            category_parts = []
                            for part in parts:
                                if any(c.isdigit() for c in part):
                                    break
                                category_parts.append(part)
                            if category_parts:
                                category_name = ' '.join(category_parts).strip()
                        
                        value = float(row['Column1'])
                        key = (slide_index, shape_index, category_name)
                        if key in translations:
                            categories.append(translations[key])
                            updates += 1
                        else:
                            categories.append(category_name)
                        values.append(value)
                    
                    chart_data = ChartData()
                    chart_data.categories = categories
                    chart_data.add_series('Column1', values)
                    chart.replace_data(chart_data)
                    
                    try:
                        os.remove(excel_path)
                    except:
                        pass
            
            except Exception as e:
                logger.error(f"Error processing chart: {str(e)}")
                continue
    
    logger.info(f"Updated {updates} chart elements")

def update_pptx(input_pptx, output_pptx, text_jsonl, tables_jsonl, charts_jsonl):
    """Update all elements in the presentation with French translations."""
    logger.info(f"Loading presentation: {input_pptx}")
    prs = Presentation(input_pptx)
    
    # Update each type of content
    update_text_blocks(prs, text_jsonl)
    update_tables(prs, tables_jsonl)
    update_charts(prs, charts_jsonl)
    
    # Save the updated presentation
    logger.info(f"Saving updated presentation to: {output_pptx}")
    prs.save(output_pptx)
    logger.info("✅ Update complete!")

if __name__ == "__main__":
    update_pptx(
        input_pptx="/Users/jasonhao/Desktop/LLM experiments/translatedocs/forst/slides/survey-phase2-eng-PPT (2).pptx",
        output_pptx="survey-phase2-fr-PPT-complete_layout.pptx",
        text_jsonl="output/merged_text_blocks_with_translations.jsonl",
        tables_jsonl="output/merged_tables_with_translations.jsonl",
        charts_jsonl="output/merged_chart_titles.jsonl"
    ) 