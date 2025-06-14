import json
import logging
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.dml.color import MSO_COLOR_TYPE
from typing import Dict, Any, List

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ContentExtractor:
    def __init__(self):
        """Initialize the content extractor."""
        pass

    def extract_text_blocks(self, pptx_path: str, output_jsonl: str):
        """Extract text blocks with formatting (from block_extract.py)."""
        prs = Presentation(pptx_path)
        blocks = []
        
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_index, run in enumerate(paragraph.runs):
                        # Skip runs that contain only spaces
                        if not run.text.strip():
                            continue
                            
                        block = {
                            "slide_index": slide_index,
                            "shape_index": shape_index,
                            "paragraph_index": paragraph_index,
                            "run_index": run_index,
                            "font": run.font.name if run.font and run.font.name else None,
                            "size": run.font.size.pt if run.font and run.font.size else None,
                            "bold": bool(run.font.bold) if run.font.bold is not None else False,
                            "italic": bool(run.font.italic) if run.font.italic is not None else False,
                            "underline": bool(run.font.underline) if run.font.underline is not None else False,
                            "color": None,  # You can add color extraction if needed
                            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                            "text": run.text,
                            "french_text": f"[FR]{run.text}" if run.text else None
                        }
                        blocks.append(block)

        with open(output_jsonl, "w", encoding="utf-8") as f:
            for block in blocks:
                f.write(json.dumps(block, ensure_ascii=False) + "\n")
        logger.info(f"✅ Saved {len(blocks)} text blocks to {output_jsonl}")

    def extract_tables(self, pptx_path: str, full_output_jsonl: str, runs_output_jsonl: str):
        """Extract tables with full formatting and runs (from Table_extract_1.py)."""
        prs = Presentation(pptx_path)
        tables = []
        runs_for_translation = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if not shape.has_table:
                    continue

                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                data = []
                
                for r in range(rows):
                    row_data = []
                    for c in range(cols):
                        cell = table.cell(r, c)
                        cell_paragraphs = []
                        if cell.text_frame:
                            for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                paragraph_id = f"s{slide_index}_sh{shape_index}_r{r}_c{c}_p{para_idx}"
                                para_info = {
                                    "paragraph_id": paragraph_id,
                                    "text": paragraph.text,
                                    "is_bullet": False,
                                    "bullet_level": paragraph.level,
                                    "runs": []
                                }
                                # Bullet detection
                                try:
                                    p = paragraph._p
                                    if p.pPr is not None:
                                        if 'buChar' in p.pPr.xml or 'buAutoNum' in p.pPr.xml or 'buBlip' in p.pPr.xml:
                                            para_info["is_bullet"] = True
                                except Exception as e:
                                    logger.error(f"Error checking XML: {e}")

                                # Extract runs
                                for run_idx, run in enumerate(paragraph.runs):
                                    run_id = f"{paragraph_id}_run{run_idx}"
                                    run_info = {
                                        "run_id": run_id,
                                        "run_idx": f"s{slide_index}_sh{shape_index}_r{r}_c{c}_p{para_idx}_run{run_idx}",
                                        "text": run.text,
                                        "french_text": f"[FR]{run.text}",
                                        "font": None,
                                        "size": None,
                                        "bold": None,
                                        "italic": None,
                                        "underline": None,
                                        "color": None
                                    }
                                    
                                    # Get run formatting
                                    font = run.font
                                    if font:
                                        run_info["font"] = font.name if font.name else None
                                        run_info["size"] = font.size.pt if font.size else None
                                        run_info["bold"] = bool(font.bold) if font.bold is not None else None
                                        run_info["italic"] = bool(font.italic) if font.italic is not None else None
                                        run_info["underline"] = bool(font.underline) if font.underline is not None else None
                                        
                                        # Get color information
                                        if font.color:
                                            if font.color.type == MSO_COLOR_TYPE.RGB and font.color.rgb:
                                                run_info["color"] = str(font.color.rgb)
                                            else:
                                                run_info["color"] = str(font.color.type)
                                    
                                    # Add hyperlink if present
                                    if run.hyperlink:
                                        run_info["url"] = run.hyperlink.address
                                    
                                    para_info["runs"].append(run_info)
                                    runs_for_translation.append({
                                        'run_idx': run_info['run_idx'],
                                        'text': run_info['text'],
                                        'french_text': run_info['french_text']
                                    })
                                
                                cell_paragraphs.append(para_info)
                        cell_info = {
                            "paragraphs": cell_paragraphs
                        }
                        row_data.append(cell_info)
                    
                    data.append(row_data)

                tables.append({
                    'slide_index': slide_index,
                    'shape_index': shape_index,
                    'rows': rows,
                    'cols': cols,
                    'data': data,
                    'french_data': []  # Empty list for French data
                })

        # Save full table structure
        with open(full_output_jsonl, 'w', encoding='utf-8') as f:
            for table in tables:
                f.write(json.dumps(table, ensure_ascii=False) + '\n')
        logger.info(f"✅ Saved {len(tables)} tables to {full_output_jsonl}")

        # Save runs for translation
        with open(runs_output_jsonl, 'w', encoding='utf-8') as f:
            for run in runs_for_translation:
                f.write(json.dumps(run, ensure_ascii=False) + '\n')
        logger.info(f"✅ Saved {len(runs_for_translation)} table runs to {runs_output_jsonl}")

    def extract_chart_titles(self, pptx_path: str, output_jsonl: str):
        """Extract chart titles (from extract_chart_titles.py)."""
        prs = Presentation(pptx_path)
        chart_titles = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if not isinstance(shape, GraphicFrame):
                    continue
                try:
                    chart = shape.chart
                except (ValueError, AttributeError):
                    continue
                if chart.has_title:
                    title = chart.chart_title.text_frame.text.strip()
                    if title:
                        chart_titles.append({
                            'slide_index': slide_index,
                            'shape_index': shape_index,
                            'element': 'chart_title',
                            'text': title,
                            'french_text': '[FR] ' + title
                        })

        with open(output_jsonl, 'w', encoding='utf-8') as f:
            for block in chart_titles:
                f.write(json.dumps(block, ensure_ascii=False) + '\n')
        logger.info(f"✅ Saved {len(chart_titles)} chart titles to {output_jsonl}")

def main():
    """Main function to demonstrate usage."""
    extractor = ContentExtractor()
    pptx_path = "slides/survey-phase2-eng-PPT (2).pptx"
    
    # Extract all content types
    extractor.extract_text_blocks(
        pptx_path=pptx_path,
        output_jsonl="input/text_blocks.jsonl"
    )
    
    extractor.extract_tables(
        pptx_path=pptx_path,
        full_output_jsonl="input/tables_full.jsonl",
        runs_output_jsonl="input/tables_runs.jsonl"
    )
    
    extractor.extract_chart_titles(
        pptx_path=pptx_path,
        output_jsonl="input/chart_titles.jsonl"
    )

if __name__ == "__main__":
    main() 