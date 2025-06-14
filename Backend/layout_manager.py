from pptx import Presentation
from pptx.util import Pt
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.chart import XL_CHART_TYPE
import json
import logging
import os

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Text Layout Functions
def estimate_text_lines(text, font_size, container_width, margin_left, margin_right):
    """Estimate number of lines based on text content and container properties."""
    if not text:
        return 1
        
    # Calculate usable width (container width minus margins)
    usable_width = container_width - (margin_left + margin_right)
    
    # Convert PowerPoint EMU to points (1 point = 12700 EMU)
    usable_width_pt = usable_width / 12700
    
    # For Arial, each character is roughly 0.8 * font_size in width
    char_width = font_size * 0.45
    
    # Calculate how many characters can fit in one line
    chars_per_line = int(usable_width_pt / char_width)
    
    # Get character count
    char_count = len(text)
    
    # Calculate number of lines needed
    if char_count > chars_per_line:
        return (char_count // chars_per_line) + 1
    return 1

def adjust_text_font_size(shape, eng_lines, fr_lines, initial_font_size):
    """Adjust font size until French text matches English line count."""
    if not shape.has_text_frame:
        return initial_font_size
        
    text_frame = shape.text_frame
    current_font_size = initial_font_size
    max_attempts = 20
    
    for _ in range(max_attempts):
        current_fr_lines = estimate_text_lines(
            text_frame.text,
            current_font_size,
            shape.width,
            text_frame.margin_left,
            text_frame.margin_right
        )
        
        if current_fr_lines == eng_lines:
            break
            
        if current_fr_lines > eng_lines:
            current_font_size -= 1
        else:
            current_font_size += 1
            
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    run.font.size = Pt(current_font_size)
    
    return current_font_size

# Table Layout Functions
def emu_to_points(emu):
    """Convert EMU (English Metric Units) to points."""
    return emu / 12700

def estimate_table_lines(text, width_pt, font_size):
    """Estimate number of lines needed for text in a table cell."""
    if not text:
        return 0
        
    char_width = 0.85 * font_size
    chars_per_line = int(width_pt / char_width)
    estimated_lines = len(text) / chars_per_line
    return int(estimated_lines) + (1 if estimated_lines % 1 > 0 else 0)

def adjust_table_font_size(eng_lines, fr_text, width_pt, template_font_size):
    """Adjust French font size to match English line count in tables."""
    if not fr_text:
        return template_font_size
        
    current_font_size = template_font_size
    min_font_size = 8
    
    while current_font_size >= min_font_size:
        fr_lines = estimate_table_lines(fr_text, width_pt, current_font_size)
        if fr_lines <= eng_lines:
            break
        current_font_size -= 1
    
    return current_font_size

def adjust_table_font_by_dimensions(eng_cell_height, fr_cell_height, fr_text, width_pt, template_font_size, eng_lines):
    """Adjust French font size based on cell height comparison."""
    if not fr_text or fr_cell_height <= eng_cell_height:
        return template_font_size
        
    current_font_size = template_font_size
    min_font_size = 8
    
    height_ratio = eng_cell_height / fr_cell_height
    
    while current_font_size >= min_font_size:
        estimated_lines = estimate_table_lines(fr_text, width_pt, current_font_size)
        if estimated_lines <= eng_lines:
            break
        current_font_size -= 1
    
    return current_font_size

# Chart Functions
def get_chart_type_name(chart_type):
    """Convert chart type number to readable name."""
    chart_types = {
        -4120: "DOUGHNUT",
        -4121: "PIE",
        -4169: "LINE",
        -4170: "COLUMN",
        -4171: "BAR",
        -4172: "AREA",
        -4173: "SCATTER",
        -4174: "STOCK",
        -4175: "RADAR",
        -4176: "SURFACE",
        -4177: "BUBBLE",
        -4178: "CANDLESTICK",
        -4179: "CYLINDER",
        -4180: "CONE",
        -4181: "PYRAMID"
    }
    return chart_types.get(chart_type, f"UNKNOWN ({chart_type})")

def is_chart_shape(shape):
    """Check if a shape is a chart."""
    if not isinstance(shape, GraphicFrame):
        return False
    try:
        return hasattr(shape, 'chart') and shape.chart is not None
    except:
        return False

# Main Processing Functions
def process_text_layout(prs, language, eng_results=None):
    """Process text layout in the presentation."""
    logger.info(f"Processing text layout for {language} presentation...")
    results = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide.slide_id // 256
        
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
                
            properties = {
                "slide_number": slide_number,
                "slide_index": slide_idx,
                "shape_index": shape_idx,
                "shape_type": str(shape.shape_type),
                "language": language,
                "paragraphs": []
            }
            
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                full_text = ""
                para_font_size = None
                runs_info = []
                
                for run in paragraph.runs:
                    run_info = {
                        "text": run.text,
                        "font_size": run.font.size.pt if run.font and run.font.size else None
                    }
                    runs_info.append(run_info)
                    
                    if run.font and run.font.size and para_font_size is None:
                        para_font_size = run.font.size.pt
                    full_text += run.text
                
                if full_text:
                    properties["paragraphs"].append({
                        "text": full_text,
                        "char_count": len(full_text),
                        "estimated_lines": estimate_text_lines(
                            full_text,
                            para_font_size or 12,
                            shape.width,
                            text_frame.margin_left,
                            text_frame.margin_right
                        ),
                        "runs": runs_info
                    })
            
            if language == "fr" and eng_results:
                key = f"{slide_idx}_{shape_idx}"
                if key in eng_results:
                    eng_shape = eng_results[key]
                    eng_text = " ".join(para["text"] for para in eng_shape["paragraphs"])
                    fr_text = " ".join(para["text"] for para in properties["paragraphs"])
                    
                    initial_font_size = 12
                    if properties["paragraphs"] and properties["paragraphs"][0]["runs"]:
                        initial_font_size = properties["paragraphs"][0]["runs"][0]["font_size"] or 12
                    
                    eng_lines = estimate_text_lines(
                        eng_text,
                        initial_font_size,
                        shape.width,
                        text_frame.margin_left,
                        text_frame.margin_right
                    )
                    fr_lines = estimate_text_lines(
                        fr_text,
                        initial_font_size,
                        shape.width,
                        text_frame.margin_left,
                        text_frame.margin_right
                    )
                    
                    if eng_lines != fr_lines:
                        new_font_size = adjust_text_font_size(
                            shape,
                            eng_lines,
                            fr_lines,
                            initial_font_size
                        )
                        
                        for para in properties["paragraphs"]:
                            para["adjusted_font_size"] = new_font_size
                            para["estimated_lines"] = eng_lines
            
            key = f"{slide_idx}_{shape_idx}"
            results[key] = properties
    
    return results

def process_table_layout(prs, language, eng_results=None):
    """Process table layout in the presentation."""
    logger.info(f"Processing table layout for {language} presentation...")
    results = []
    
    for slide_idx, slide in enumerate(prs.slides):
        tables = [shape for shape in slide.shapes if shape.has_table]
        
        for table_idx, table in enumerate(tables):
            table_info = {
                "slide": slide_idx + 1,
                "table": table_idx + 1,
                "language": language,
                "rows": len(table.table.rows),
                "columns": len(table.table.columns),
                "cells": []
            }
            
            column_widths = [col.width for col in list(table.table.columns)]
            
            # Find corresponding English table if available
            eng_table = None
            if language == "fr" and eng_results:
                eng_table = next(
                    (t for t in eng_results if t["slide"] == slide_idx + 1 and t["table"] == table_idx + 1),
                    None
                )
            
            for row_idx, row in enumerate(table.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    text_frame = cell.text_frame
                    original_font_size = 12
                    
                    if text_frame:
                        for paragraph in text_frame.paragraphs:
                            if paragraph.text and paragraph.runs and paragraph.runs[0].font.size:
                                original_font_size = paragraph.runs[0].font.size.pt
                                break
                    
                    cell_text = cell.text if cell.text else ""
                    width_pt = round(emu_to_points(column_widths[col_idx]), 2)
                    height_pt = round(emu_to_points(row.height), 2)
                    
                    if language == "fr" and eng_table:
                        # Find corresponding English cell
                        eng_cell = next(
                            (c for c in eng_table["cells"] if c["row"] == row_idx and c["column"] == col_idx),
                            None
                        )
                        
                        if eng_cell and eng_cell["text"]:
                            template_font_size = eng_cell["original_font_size"]
                            dimension_adjusted_font_size = adjust_table_font_by_dimensions(
                                eng_cell["height_pt"],
                                height_pt,
                                cell_text,
                                width_pt,
                                template_font_size,
                                eng_cell["original_lines"]
                            )
                            
                            line_adjusted_font_size = adjust_table_font_size(
                                eng_cell["original_lines"],
                                cell_text,
                                width_pt,
                                template_font_size
                            )
                            
                            adjusted_font_size = max(dimension_adjusted_font_size, line_adjusted_font_size)
                            adjusted_lines = estimate_table_lines(cell_text, width_pt, adjusted_font_size)
                            
                            while adjusted_lines > eng_cell["original_lines"] and adjusted_font_size > 8:
                                adjusted_font_size -= 1
                                adjusted_lines = estimate_table_lines(cell_text, width_pt, adjusted_font_size)
                            
                            original_lines = estimate_table_lines(cell_text, width_pt, original_font_size)
                        else:
                            original_lines = estimate_table_lines(cell_text, width_pt, original_font_size)
                            adjusted_lines = original_lines
                            adjusted_font_size = original_font_size
                    else:
                        original_lines = estimate_table_lines(cell_text, width_pt, original_font_size)
                        adjusted_lines = original_lines
                        adjusted_font_size = original_font_size
                    
                    cell_info = {
                        "row": row_idx,
                        "column": col_idx,
                        "width_pt": width_pt,
                        "height_pt": height_pt,
                        "text": cell_text,
                        "original_font_size": original_font_size,
                        "adjusted_font_size": adjusted_font_size,
                        "original_lines": original_lines,
                        "adjusted_lines": adjusted_lines
                    }
                    table_info["cells"].append(cell_info)
            
            results.append(table_info)
    
    return results

def process_chart_layout(prs):
    """Process chart layout in the presentation."""
    logger.info("Processing chart layout...")
    chart_details = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not is_chart_shape(shape):
                continue
                
            try:
                chart = shape.chart
                chart_info = {
                    "slide_number": slide_idx + 1,
                    "shape_number": shape_idx + 1,
                    "chart_type": get_chart_type_name(chart.chart_type),
                    "has_title": chart.has_title,
                    "title": chart.chart_title.text_frame.text if chart.has_title else None,
                    "has_legend": chart.has_legend,
                    "plot_area": {},
                    "series": [],
                    "axes": {}
                }
                
                if hasattr(chart, 'plot_area'):
                    plot = chart.plot_area
                    chart_info["plot_area"] = {
                        "has_data_table": getattr(plot, 'has_data_table', False),
                        "has_data_labels": getattr(plot, 'has_data_labels', False),
                        "has_legend": getattr(plot, 'has_legend', False),
                        "has_major_gridlines": getattr(plot, 'has_major_gridlines', False),
                        "has_minor_gridlines": getattr(plot, 'has_minor_gridlines', False)
                    }
                
                if hasattr(chart, 'series'):
                    for series in chart.series:
                        series_info = {
                            "name": series.name,
                            "has_data_labels": getattr(series, 'has_data_labels', False),
                            "marker_style": str(series.marker_style) if hasattr(series, 'marker_style') else None,
                            "values": []
                        }
                        
                        if hasattr(series, 'values'):
                            try:
                                series_info["values"] = [str(val) for val in series.values]
                            except:
                                series_info["values"] = ["Unable to extract values"]
                        
                        chart_info["series"].append(series_info)
                
                for axis_name in ['category_axis', 'value_axis', 'second_category_axis', 'second_value_axis']:
                    if hasattr(chart, axis_name):
                        axis = getattr(chart, axis_name)
                        axis_info = {
                            "has_title": getattr(axis, 'has_title', False),
                            "title": axis.axis_title.text_frame.text if getattr(axis, 'has_title', False) else None,
                            "has_major_gridlines": getattr(axis, 'has_major_gridlines', False),
                            "has_minor_gridlines": getattr(axis, 'has_minor_gridlines', False),
                            "format": str(axis.format) if hasattr(axis, 'format') else None
                        }
                        chart_info["axes"][axis_name] = axis_info
                
                try:
                    chart_data = chart.part.chart_workbook.xlsx_part.blob
                    chart_info["has_data"] = True
                except:
                    chart_info["has_data"] = False
                
                chart_details.append(chart_info)
                logger.info(f"Processed chart on slide {slide_idx + 1}, shape {shape_idx + 1}")
                
            except Exception as e:
                logger.error(f"Error processing chart on slide {slide_idx + 1}, shape {shape_idx + 1}: {str(e)}")
                continue
    
    return chart_details

def apply_layout_adjustments(fr_prs, text_results, table_results):
    """Apply layout adjustments to the French presentation."""
    logger.info("Applying layout adjustments...")
    
    # Apply text adjustments
    for key, shape_info in text_results.items():
        slide_idx, shape_idx = map(int, key.split('_'))
        shape = fr_prs.slides[slide_idx].shapes[shape_idx]
        
        if shape.has_text_frame:
            for para_idx, para_info in enumerate(shape_info["paragraphs"]):
                if "adjusted_font_size" in para_info:
                    for run in shape.text_frame.paragraphs[para_idx].runs:
                        if run.font:
                            run.font.size = Pt(para_info["adjusted_font_size"])
    
    # Apply table adjustments
    for table_info in table_results:
        if table_info["language"] != "fr":
            continue
            
        slide_idx = table_info["slide"] - 1
        table_idx = table_info["table"] - 1
        
        try:
            table = fr_prs.slides[slide_idx].shapes[table_idx].table
            for cell_info in table_info["cells"]:
                if cell_info["adjusted_font_size"] != cell_info["original_font_size"]:
                    cell = table.cell(cell_info["row"], cell_info["column"])
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font:
                                run.font.size = Pt(cell_info["adjusted_font_size"])
        except Exception as e:
            logger.error(f"Error applying table adjustments: {str(e)}")
            continue

def main(
    input_pptx: str,
    output_pptx: str,
    layout_dir: str = "layout"
):
    """
    Main function to process and adjust layouts.
    
    Args:
        input_pptx (str): Path to input English PowerPoint file
        output_pptx (str): Path to output French PowerPoint file
        layout_dir (str): Directory to save layout information
    """
    # Create layout directory if it doesn't exist
    os.makedirs(layout_dir, exist_ok=True)
    
    # Process English presentation
    eng_prs = Presentation(input_pptx)
    eng_text_results = process_text_layout(eng_prs, "en")
    eng_table_results = process_table_layout(eng_prs, "en")
    eng_chart_details = process_chart_layout(eng_prs)
    
    # Process French presentation
    fr_prs = Presentation(output_pptx)
    fr_text_results = process_text_layout(fr_prs, "fr", eng_text_results)
    fr_table_results = process_table_layout(fr_prs, "fr", eng_table_results)
    fr_chart_details = process_chart_layout(fr_prs)
    
    # Save results
    with open(os.path.join(layout_dir, "text_layout.jsonl"), "w", encoding="utf-8") as f:
        for key in eng_text_results:
            if key in fr_text_results:
                f.write(json.dumps(eng_text_results[key], ensure_ascii=False) + "\n")
                f.write(json.dumps(fr_text_results[key], ensure_ascii=False) + "\n")
    
    with open(os.path.join(layout_dir, "table_layout.jsonl"), "w", encoding="utf-8") as f:
        for eng_table, fr_table in zip(eng_table_results, fr_table_results):
            f.write(json.dumps(eng_table, ensure_ascii=False) + "\n")
            f.write(json.dumps(fr_table, ensure_ascii=False) + "\n")
    
    with open(os.path.join(layout_dir, "chart_layout.json"), "w", encoding="utf-8") as f:
        json.dump({
            "english": eng_chart_details,
            "french": fr_chart_details
        }, f, indent=2, ensure_ascii=False)
    
    # Apply adjustments to French presentation
    apply_layout_adjustments(fr_prs, fr_text_results, fr_table_results)
    
    # Save adjusted French presentation
    layout_output_pptx = output_pptx.replace(".pptx", "_layout.pptx")
    fr_prs.save(layout_output_pptx)
    logger.info(f"✅ Layout processing complete! Adjusted file saved to {layout_output_pptx}")
    
    return layout_output_pptx

if __name__ == "__main__":
    # Example usage
    input_pptx = "slides/PPT-3-Government-in-Canada1 (1).pptx"
    output_pptx = "PPT-3-Government-in-Canada1 (1)_fr.pptx"
    main(input_pptx, output_pptx) 